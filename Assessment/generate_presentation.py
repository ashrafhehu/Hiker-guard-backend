import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def build_presentation():
    input_file = "BC1-D1 (Certified Assessment 1) Template_v1.1.pptx"
    output_file = "AGAIF2026_BC1_CA1_MY-411_Muhammad Ashraf.pptx"
    
    prs = pptx.Presentation(input_file)
    slide = prs.slides[0]

    # Harmonized Palette
    LABEL_COLOR = RGBColor(30, 58, 138)      # Deep Navy Blue #1e3a8a
    BODY_COLOR = RGBColor(30, 41, 59)       # Dark Slate / Charcoal #1e293b
    REF_COLOR = RGBColor(71, 85, 105)       # Slate Gray #475569

    # Clean Unicode bullet
    BULLET = "\u2022 "

    # Content data for each dimension
    content_map = {
        5: [ # Dimension 1: Problem Statement & Geographic Context
            ("Real-World Problem: ", "Mountainous terrain, deep valleys, and dense rainforest canopies severely degrade or sever cellular connectivity along remote hiking corridors (e.g., Gunung Tahan, Bukit Larut). Hikers losing signal without warning cannot transmit live GPS telemetry, request emergency aid, or access online safety maps."),
            ("Commercial Data Deficiency: ", "Commercial MNO coverage maps are macro-scale, heavily aggregated, and overly optimistic in rural/montane areas, creating a false sense of coverage security."),
            ("Geographic Context: ", "Tropical montane & rainforest hiking trails in Malaysia characterized by steep elevation gradients, severe ridge/valley line-of-sight obstruction, and dense environmental vegetation clutter.")
        ],
        14: [ # Dimension 2: Stakeholders & Users
            ("Outdoor Hikers & Guides: ", "Receive pre-hike connectivity gap alerts, offline map caching prompts, and automatic trajectory queuing before entering dead zones."),
            ("Park & Forestry Authorities: ", "Utilize trail-segment risk maps for visitor management, ranger safety, and trail infrastructure planning."),
            ("Search & Rescue (SAR) Teams: ", "Access last-known connectivity boundary layers to narrow search radii during emergency search operations (BOMBA, SMART, APM)."),
            ("Telecom Regulators & MNOs: ", "Utilize empirical gap evidence for targeted rural coverage infrastructure planning.")
        ],
        23: [ # Dimension 3: Spatial Context & Scale
            ("Spatial Resolution: ", "Resampled ~250m deterministic vector trail segments (segment_id) derived from high-precision GPX trail geometries."),
            ("Topographic & Environmental Variables: ", "Elevation, slope, Terrain Ruggedness Index (TRI), line-of-sight terrain obstruction profiles to cell towers, and forest canopy clutter."),
            ("Geographic Domain Transfer: ", "Bridges data-rich source geographies (Anatel 4G Brazil, FCC BDC US 4G LTE, Ofcom UK 4G) to data-sparse target tropical mountain trails in Malaysia via domain adaptation.")
        ],
        32: [ # Dimension 4: Data Requirements
            ("Copernicus 30m DEM: ", "Elevation, slope, TRI, line-of-sight elevation obstruction profiles."),
            ("ESA WorldCover 10m: ", "Land-cover classification, tree cover %, environmental clutter proxy."),
            ("OpenCellID: ", "Distance to nearest cell tower, tower density (1km/3km/5km radii), radio tech."),
            ("Ookla Mobile Performance: ", "Speed, latency, device count (positive coverage evidence only; missing = unobserved)."),
            ("GPX Geometries: ", "Resampled ~250m trail linestring segments."),
            ("Cross-Country Labels: ", "Anatel 4G Brazil (primary weak label), FCC BDC US 4G LTE (secondary weak label), Ofcom UK 4G (benchmark validation).")
        ],
        41: [ # Dimension 5: GeoAI Method & Analytical Approach
            ("Model Architecture: ", "Gradient-Boosted Decision Trees (LightGBM/XGBoost) trained via domain-adversarial transfer learning across source-target geographic features."),
            ("Spatial Feature Engineering: ", "Topographic line-of-sight obstacle index (elevation profile between segment centroid and cell tower), forest clutter density, tower distance decay, and Ookla signal presence flags."),
            ("Validation & Domain Adaptation: ", "Spatial cross-validation (GroupKFold by mountain block), out-of-distribution (OOD) distance metrics, and domain similarity checks.")
        ],
        50: [ # Dimension 6: GeoAI Solution Design & Decision Support
            ("Cautious 3-Tier Classification: ", "Outputs likely_covered (risk <= 0.35 + evidence), uncertain (sparse/conflicting/high-OOD data), and predicted_gap (risk >= 0.70 + evidence)."),
            ("Actionable Smartphone Warnings: ", "Pre-hike alert: \"Predicted gap in 600m: Download offline map, cache position, and check battery.\""),
            ("Offline Trajectory Logging: ", "Automatic background location queuing with auto-sync upon reconnection."),
            ("Park Management Dashboard: ", "Spatial gap heatmaps & trail audit reports for SAR support.")
        ],
        59: [ # Dimension 7: Technology Stack
            ("Data Pipeline: ", "Python, PyGIS (GeoPandas, Shapely, Rasterio), Copernicus DEM & ESA WorldCover spatial ingestion pipelines."),
            ("GeoAI Modeling & Registry: ", "LightGBM, XGBoost, Scikit-Learn, MLflow model registry, Pydantic data contract validation."),
            ("Inference & App Delivery: ", "FastAPI REST microservice (exposing versioned prediction contracts), Mobile Client (Flutter/Native) with local SpatiaLite/SQLite offline cache, background GPS tracker, and auto-sync queue.")
        ],
        68: [ # Dimension 8: Ethics, Privacy & Sustainability
            ("Ethical & Safety Safeguards: ", "Cautious 3-tier terminology (likely_covered, uncertain, predicted_gap). NEVER claims \"guaranteed zero coverage\" to prevent reckless hiker behavior."),
            ("Privacy & Data Governance: ", "On-device offline trajectory logging; location telemetry is anonymized and transmitted only upon explicit user consent during online sync."),
            ("Sustainability & Cost Feasibility: ", "Phone-only solution requiring zero physical trail hardware or LoRa towers. Low-carbon compute footprint using optimized tree ML inference.")
        ],
        77: [ # Dimension 9: Value & Impact
            ("Hiker Safety Improvement: ", "Prevents disorientation and emergency delays by prompting proactive map caching and family location sharing prior to gap entry."),
            ("SAR Operational Efficiency: ", "Narrows SAR search grid from hundreds of sq km down to the precise last-connected trail segment boundary, drastically cutting search time."),
            ("Feasibility & Scalability: ", "Highly feasible (100% phone-only, zero hardware), globally scalable to any hiking trail network using public open geospatial datasets.")
        ],
        86: [ # Dimension 10: Implementation Roadmap
            ("Phase 1: MVP Baseline (M1-M3): ", "Open dataset pipeline, cross-country transfer model training, evaluation on 4 Peninsular Malaysia benchmark routes (Gunung Tahan, Bukit Larut, Kledang, Bukit Kerinchi), mobile offline alert client."),
            ("Phase 2: Field Calibration & Scaling (M4-M6): ", "Hiker crowdsourced field verification, MNO/authority data sharing (MCMC/JENDELA integration), expanding coverage to 50+ national park trails across Malaysia."),
            ("Phase 3: Operational SAR & Global Rollout (M7-M12): ", "Integration with national park SAR management platforms, expanding to international hiking corridors, predictive SAR risk mapping.")
        ],
        108: [ # Dimension 11: Final Pitch
            ("Elevator Pitch: ", "JEJAK is a phone-only GeoAI decision-support platform that transforms terrain, land cover, and cell infrastructure telemetry into predictive cellular gap warnings and offline safety guidance for hikers. By leveraging cross-country transfer learning and global open geospatial data, JEJAK protects lives in remote mountain trails without requiring costly trail hardware or over-promising coverage certainty."),
            ("References & Data Sources: ", "Copernicus DEM 30m, ESA WorldCover 10m, OpenCellID, Ookla Mobile Performance Tiles, Anatel Brazil 4G, FCC BDC US 4G LTE, Ofcom UK 4G Measurements.")
        ]
    }

    # 1. Update Header Metadata Box (TextBox 122)
    for s in slide.shapes:
        if s.shape_id == 122 and s.has_text_frame:
            tf = s.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = Pt(6)
            tf.margin_right = Pt(6)
            tf.margin_top = Pt(4)
            tf.margin_bottom = Pt(4)
            
            p1 = tf.paragraphs[0]
            p1.alignment = PP_ALIGN.LEFT
            r1_label = p1.add_run()
            r1_label.text = "Project Title: "
            r1_label.font.bold = True
            r1_label.font.name = "Arial"
            r1_label.font.size = Pt(8.5)
            r1_label.font.color.rgb = LABEL_COLOR
            
            r1_val = p1.add_run()
            r1_val.text = "JEJAK – Phone-Only GeoAI Connectivity Intelligence for Hiking Safety"
            r1_val.font.bold = False
            r1_val.font.name = "Arial"
            r1_val.font.size = Pt(8.5)
            r1_val.font.color.rgb = BODY_COLOR

            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            p2.space_before = Pt(2)
            r2_label = p2.add_run()
            r2_label.text = "Name: "
            r2_label.font.bold = True
            r2_label.font.name = "Arial"
            r2_label.font.size = Pt(8.5)
            r2_label.font.color.rgb = LABEL_COLOR
            
            r2_val = p2.add_run()
            r2_val.text = "Muhammad Ashraf"
            r2_val.font.bold = False
            r2_val.font.name = "Arial"
            r2_val.font.size = Pt(8.5)
            r2_val.font.color.rgb = BODY_COLOR

            p3 = tf.add_paragraph()
            p3.alignment = PP_ALIGN.LEFT
            p3.space_before = Pt(2)
            r3_label = p3.add_run()
            r3_label.text = "AGAIF Ref. Code: "
            r3_label.font.bold = True
            r3_label.font.name = "Arial"
            r3_label.font.size = Pt(8.5)
            r3_label.font.color.rgb = LABEL_COLOR
            
            r3_val = p3.add_run()
            r3_val.text = "MY-411"
            r3_val.font.bold = False
            r3_val.font.name = "Arial"
            r3_val.font.size = Pt(8.5)
            r3_val.font.color.rgb = BODY_COLOR

    # 2. Populate Card Content Text Boxes
    for s in slide.shapes:
        if s.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
            for sub in s.shapes:
                if sub.shape_id in content_map and sub.has_text_frame:
                    if sub.shape_id == 108:
                        # Box 11: Final Pitch - Position text frame to the right of '11 FINAL PITCH' header
                        sub.left = 1250000
                        sub.top = 30000
                        sub.width = 3450000
                        sub.height = 800000
                        font_size = Pt(8.5)
                    else:
                        # Cards 1 to 10: Position text frame below the top colored header banner!
                        sub.top = 220000
                        sub.height = 800000
                        if sub.shape_id == 32:  # Data requirements (6 bullets)
                            font_size = Pt(8)
                        else:
                            font_size = Pt(8.5)
                        
                    bullets = content_map[sub.shape_id]
                    tf = sub.text_frame
                    tf.clear()
                    tf.word_wrap = True
                    tf.vertical_anchor = MSO_ANCHOR.TOP  # Top align inside the card body!
                    tf.margin_left = Pt(6)
                    tf.margin_right = Pt(6)
                    tf.margin_top = Pt(4)
                    tf.margin_bottom = Pt(4)
                    
                    for idx, (prefix, text) in enumerate(bullets):
                        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                        p.alignment = PP_ALIGN.LEFT  # Strict Left Alignment!
                        p.space_after = Pt(2.5)
                        p.line_spacing = 1.12
                        
                        r_pre = p.add_run()
                        r_pre.text = BULLET + prefix
                        r_pre.font.name = "Arial"
                        r_pre.font.bold = True
                        r_pre.font.size = font_size
                        if sub.shape_id == 108 and "References" in prefix:
                            r_pre.font.color.rgb = REF_COLOR
                            r_pre.font.italic = True
                        else:
                            r_pre.font.color.rgb = LABEL_COLOR
                            
                        r_txt = p.add_run()
                        r_txt.text = text
                        r_txt.font.name = "Arial"
                        r_txt.font.bold = False
                        r_txt.font.size = font_size
                        if sub.shape_id == 108 and "References" in prefix:
                            r_txt.font.color.rgb = REF_COLOR
                            r_txt.font.italic = True
                        else:
                            r_txt.font.color.rgb = BODY_COLOR

    prs.save(output_file)
    print(f"Successfully generated clean presentation: {output_file}")

if __name__ == "__main__":
    build_presentation()
